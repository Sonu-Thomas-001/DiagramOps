import io
import base64
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

async def create_ppt_from_image(image_data: str, title: str = "System Architecture") -> io.BytesIO:
    """
    Generates a PowerPoint presentation with a single slide containing the title and the centered diagram.
    Handles both HTTP URLs and Base64 Data URIs.
    """
    prs = Presentation()
    
    # Set widescreen 16:9 aspect ratio (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Use a blank slide layout (usually index 6)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # --- Add Title ---
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59) # Slate-800
    
    # --- Process Image Data ---
    image_stream = io.BytesIO()
    
    if image_data.startswith("http://") or image_data.startswith("https://"):
        # Download image from URL
        async with httpx.AsyncClient() as client:
            response = await client.get(image_data)
            response.raise_for_status()
            image_stream.write(response.content)
    elif image_data.startswith("data:image"):
        # Decode Base64 Data URI
        header, encoded = image_data.split(",", 1)
        image_stream.write(base64.b64decode(encoded))
    else:
        # Assume raw base64
        image_stream.write(base64.b64decode(image_data))
        
    image_stream.seek(0)
    
    # --- Add Centered Image ---
    # Target image size: 11 inches wide, maintaining aspect ratio
    # Center horizontally: (13.33 - 11) / 2 = 1.165 inches
    # Position vertically below title: ~1.5 inches from top
    
    try:
        slide.shapes.add_picture(
            image_stream, 
            left=Inches(1.165), 
            top=Inches(1.5), 
            width=Inches(11)
        )
    except Exception as e:
        raise ValueError(f"Failed to add image to PPT: {str(e)}")
    
    # --- Save to BytesIO ---
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

# --- Example Usage ---
if __name__ == "__main__":
    import asyncio
    
    async def test_ppt():
        # Using a placeholder image URL for testing
        test_url = "https://picsum.photos/seed/diagram/1024/576"
        ppt_bytes = await create_ppt_from_image(test_url, "E-Commerce Microservices")
        
        with open("test_architecture.pptx", "wb") as f:
            f.write(ppt_bytes.read())
        print("PPT generated successfully: test_architecture.pptx")
        
    asyncio.run(test_ppt())
